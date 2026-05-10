"""Phase 2: multimodal ingest.

- **PDF** (`pypdf`): extract text per-page, treat each page as a `_RawDoc` with
  `modtype="pdf_page"`. Works offline.
- **Slides (.pptx)**: optional via `python-pptx`; falls back to PDF if you
  pre-export.
- **Audio / video transcript** (`faster-whisper` or external Whisper service):
  *interface only* — `transcribe_audio()` raises if the optional dep is missing
  with a one-line install hint. Once transcribed, ingest the transcript through
  the same path as PDFs.

These functions are deliberately small and stateless so they can be called
from a CLI batch job, a Moodle event callback, or a scheduled scan of
`moodle/{moodledata}/filedir/*` (for power users who want full-corpus
coverage without going through Moodle web services).
"""

from __future__ import annotations

import logging
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from qdrant_client.http import models as qm

from .chunker import semantic_chunks
from .lineage import lineage_payload, new_ingest_run_id
from .ollama import OllamaClient
from .vectorstore import VectorStore

log = logging.getLogger(__name__)


@dataclass
class MultimodalDoc:
    course_id: int
    source_kind: str  # "pdf" | "pptx" | "audio" | "video"
    source_path: str
    title: str
    page_or_segment: str
    text: str


def extract_pdf(path: Path, *, course_id: int, title: str | None = None) -> list[MultimodalDoc]:
    try:
        import pypdf
    except ImportError as exc:
        raise RuntimeError(
            "PDF ingest requires `pypdf`. Add `pypdf` to deps and rebuild the image."
        ) from exc

    reader = pypdf.PdfReader(str(path))
    out: list[MultimodalDoc] = []
    for i, page in enumerate(reader.pages):
        text = (page.extract_text() or "").strip()
        if not text:
            continue
        out.append(
            MultimodalDoc(
                course_id=course_id,
                source_kind="pdf",
                source_path=str(path),
                title=title or path.name,
                page_or_segment=f"page {i + 1}",
                text=text,
            )
        )
    return out


def extract_pptx(path: Path, *, course_id: int, title: str | None = None) -> list[MultimodalDoc]:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as exc:
        raise RuntimeError(
            "PPTX ingest requires `python-pptx`. Install it or pre-export to PDF."
        ) from exc
    prs = Presentation(str(path))
    out: list[MultimodalDoc] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            tf = getattr(shape, "text_frame", None)
            if tf is None:
                continue
            for para in tf.paragraphs:
                t = "".join(r.text for r in para.runs).strip()
                if t:
                    parts.append(t)
        text = "\n".join(parts).strip()
        if not text:
            continue
        out.append(
            MultimodalDoc(
                course_id=course_id,
                source_kind="pptx",
                source_path=str(path),
                title=title or path.name,
                page_or_segment=f"slide {i}",
                text=text,
            )
        )
    return out


def transcribe_audio(path: Path, *, course_id: int, title: str | None = None) -> list[MultimodalDoc]:
    """Optional: requires `faster-whisper`. Otherwise raises with a hint."""
    try:
        from faster_whisper import WhisperModel  # type: ignore
    except ImportError as exc:
        raise RuntimeError(
            "Audio/video ingest requires `faster-whisper`. "
            "Install with: pip install faster-whisper"
        ) from exc
    model = WhisperModel("base", device="auto")
    segments, _info = model.transcribe(str(path))
    out: list[MultimodalDoc] = []
    for s in segments:
        text = (s.text or "").strip()
        if not text:
            continue
        out.append(
            MultimodalDoc(
                course_id=course_id,
                source_kind="audio",
                source_path=str(path),
                title=title or path.name,
                page_or_segment=f"{s.start:.1f}-{s.end:.1f}s",
                text=text,
            )
        )
    return out


def ingest_multimodal(
    docs: Iterable[MultimodalDoc],
    *,
    ollama: OllamaClient,
    store: VectorStore,
    max_chars: int = 1200,
) -> dict:
    docs_list = list(docs)
    if not docs_list:
        return {"docs": 0, "vectors": 0}
    run_id = new_ingest_run_id()
    embed_model = ollama.embed_model
    points_in: list[tuple[MultimodalDoc, int, str, str]] = []
    for d in docs_list:
        for i, c in enumerate(semantic_chunks(d.text, max_chars=max_chars, default_heading=f"{d.title} – {d.page_or_segment}")):
            points_in.append((d, i, c.heading_path, c.text))
    if not points_in:
        return {"docs": len(docs_list), "vectors": 0, "ingest_run_id": run_id}
    vecs = ollama.embed_batch([p[3] for p in points_in])
    store.ensure_collection(vector_size=len(vecs[0]))
    qpoints = []
    for (d, i, h, t), v in zip(points_in, vecs):
        import hashlib

        pid = int(hashlib.sha1(f"{d.course_id}|{d.source_kind}|{d.source_path}|{i}".encode()).hexdigest()[:15], 16)
        qpoints.append(
            qm.PointStruct(
                id=pid,
                vector=v,
                payload={
                    "course_id": d.course_id,
                    "modtype": d.source_kind,
                    "title": f"{d.title} – {d.page_or_segment}",
                    "module_name": d.title,
                    "section_name": d.source_kind,
                    "heading_path": h,
                    "url": d.source_path,
                    "chunk_no": i,
                    "text": t,
                    **lineage_payload(ingest_run_id=run_id, chunk_text=t, embed_model=embed_model),
                },
            )
        )
    store.upsert(qpoints)
    return {"docs": len(docs_list), "vectors": len(qpoints), "ingest_run_id": run_id, "embed_model": embed_model}
